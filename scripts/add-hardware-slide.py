"""提出デッキに「家に置く機器」のスライドを1枚挿入する。

既存デッキ（Berry & Cream・Georgia見出し / Meiryo本文）の意匠に合わせて
図形を素で組み立てる。挿入位置は「構成」の直前（10枚目）で、
以降のスライドに焼き込まれているページ番号を振り直す。
"""

import copy
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DECK = "docs/mimamoritai-deck.pptx"
IMG = "docs/images/hardware-slide.png"
INSERT_AT = 9  # 0-based。既存の「構成」の直前

INK = RGBColor(0x2B, 0x2B, 0x2B)
BERRY = RGBColor(0x6D, 0x2E, 0x46)
ROSE = RGBColor(0xA2, 0x67, 0x69)
MUTED = RGBColor(0x7A, 0x7A, 0x7A)
FAINT = RGBColor(0xB9, 0xB9, 0xB9)
CREAM = RGBColor(0xEC, 0xE2, 0xD0)
HAIR = RGBColor(0xE3, 0xE3, 0xE3)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (s, size, bold, font, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return box


def block(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def build(prs):
    layout = next((l for l in prs.slide_layouts if l.name == "DEFAULT"), None)
    if layout is None:
        layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)

    text(s, 0.85, 0.46, 11.63, 0.72,
         [("家に置くのは、コンセント1個だけ", 31, True, "Georgia", INK)])
    text(s, 0.85, 1.18, 11.63, 0.34,
         [("カメラもセンサーも増やさない。工事も配線もない", 13.5, False, "Meiryo", ROSE)])

    s.shapes.add_picture(IMG, Inches(0.85), Inches(1.72), width=Inches(11.63))

    # 取れるもの（4枚のカード）
    labels = [
        ("消費電力", "動いているか／止まっているか"),
        ("その日の積算電力量", "いつもの一日と重ねて比べる"),
        ("ON / OFF の時刻", "起きた時間・最後に動いた時間"),
        ("遠隔で ON / OFF", "「扇風機をつけて」で操作できる"),
    ]
    x, w, gap = 0.85, 2.83, 0.10
    for i, (head, sub) in enumerate(labels):
        cx = x + i * (w + gap)
        accent = BERRY if i == 3 else ROSE
        block(s, cx, 5.05, w, 0.95, RGBColor(0xFA, 0xF8, 0xF5), HAIR)
        block(s, cx, 5.05, w, 0.05, accent)
        text(s, cx + 0.22, 5.28, w - 0.4, 0.28,
             [(head, 13, True, "Meiryo", BERRY)])
        text(s, cx + 0.22, 5.60, w - 0.4, 0.28,
             [(sub, 10.5, False, "Meiryo", MUTED)])

    # キッカー帯
    block(s, 0.85, 6.20, 11.63, 0.52, CREAM)
    block(s, 0.85, 6.20, 0.06, 0.52, BERRY)
    kicker = text(s, 1.15, 6.32, 11.08, 0.3,
                  [("実機は SwitchBot プラグミニ (JP) 1台。ハブ経由の赤外線リモコン家電にも対応済み（本作の実機構成には含まない）",
                    12.5, True, "Meiryo", BERRY)])
    kicker.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    block(s, 0.85, 6.90, 11.63, 0.012, HAIR)
    text(s, 0.85, 6.98, 5.4, 0.3,
         [("見守り隊 / CareRoute AI", 9.5, False, "Meiryo", FAINT)])
    text(s, 11.48, 6.98, 1.0, 0.3,
         [("10", 9.5, False, "Meiryo", FAINT)], align=PP_ALIGN.RIGHT)
    return s


def move(prs, frm, to):
    ids = prs.slides._sldIdLst
    el = list(ids)[frm]
    ids.remove(el)
    ids.insert(to, el)


def renumber(prs):
    """右下に焼き込まれたページ番号を、実際の並び順に合わせる。"""
    fixed = 0
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or sh.left is None:
                continue
            if sh.left < Inches(11.2) or sh.top < Inches(6.5):
                continue
            if sh.text_frame.text.strip().isdigit():
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip().isdigit():
                            if r.text != str(i):
                                fixed += 1
                            r.text = str(i)
        # 目視用
    return fixed


def main():
    prs = Presentation(DECK)
    before = len(prs.slides._sldIdLst)
    build(prs)
    move(prs, before, INSERT_AT)
    fixed = renumber(prs)
    prs.save(DECK)
    print(f"slides={len(prs.slides._sldIdLst)} renumbered={fixed}")


if __name__ == "__main__":
    sys.exit(main())
