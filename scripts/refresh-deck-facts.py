"""提出デッキの古くなった数字とリンクを直す。

2枚目と18枚目が、いまの実物と食い違っていた。

- テスト件数が 1,004 のまま（実際は 1,065）
- コードの公開先が凍結済みのリポジトリを指していた
  （都知事杯の審査対象は mimamoritai-tokyo-opendata のほう）
- デモ動画のリンクが、先に上げた別の回（音声操作と安全確認・3分11秒）を
  「通しのデモ動画」として指していた。通しの本編は 5分51秒 で、
  いま公開リポジトリに同梱してあるファイルのほうが本物
"""

import sys
from pptx import Presentation

DECK = "docs/mimamoritai-deck.pptx"

REPO = "https://github.com/monoqlo78/mimamoritai-tokyo-opendata"
MOVIE = "docs/demo/mimamoritai-demo.mp4"

# 長いものから順に当てる。短い側が先に食ってしまうのを防ぐため。
RULES = [
    ("docs/demo/mimamoritai-demo.mp4 （字幕・ナレーション入り / 3分11秒）",
     f"リポジトリ内 {MOVIE}（字幕・ナレーション入り / 5分51秒）"),
    ("https://youtu.be/TnM-RHFZ_Lc （字幕・ナレーション入り / 3分11秒）",
     f"リポジトリ内 {MOVIE}（字幕・ナレーション入り / 5分51秒）"),
    ("通しのデモ動画（3分11秒・字幕＋ナレーション入り）",
     "通しのデモ動画（5分51秒・字幕＋ナレーション入り）"),
    ("本番環境を実際に操作して録画したもの（編集で切っていない）",
     "公開リポジトリに同梱。本番環境を実際に操作して録画したもの"),
    ("https://github.com/monoqlo78/mimamoritai-careroute-ai", REPO),
    ("https://youtu.be/TnM-RHFZ_Lc", MOVIE),
    ("1,004", "1,065"),
]


def drop_link(run):
    """文字を差し替えた場所に、古い飛び先が残らないようにする。"""
    rPr = run._r.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
    if rPr is None:
        return
    for tag in ("hlinkClick", "hlinkMouseOver"):
        el = rPr.find(
            f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}")
        if el is not None:
            rPr.remove(el)


def paragraphs(shape):
    if shape.has_text_frame:
        yield from shape.text_frame.paragraphs
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                yield from cell.text_frame.paragraphs
    if shape.shape_type == 6:  # グループ
        for sub in shape.shapes:
            yield from paragraphs(sub)


def apply(prs):
    hits = []
    for n, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            for p in paragraphs(shape):
                runs = p.runs
                if not runs:
                    continue

                # まずは1つの run で完結するものを、書式を保ったまま直す。
                for old, new in RULES:
                    for r in runs:
                        if old in r.text:
                            r.text = r.text.replace(old, new)
                            drop_link(r)
                            hits.append((n, old))

                # run がまたがって切れている場合だけ、段落ごと組み直す。
                joined = "".join(r.text for r in runs)
                for old, new in RULES:
                    if old not in joined:
                        continue
                    joined = joined.replace(old, new)
                    hits.append((n, old))
                    runs[0].text = joined
                    drop_link(runs[0])
                    for r in runs[1:]:
                        r.text = ""
                    break
    return hits


def main():
    prs = Presentation(DECK)
    hits = apply(prs)
    if not hits:
        print("no change")
        return 1
    prs.save(DECK)
    for n, old in hits:
        print(f"slide {n}: {old[:46]}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
