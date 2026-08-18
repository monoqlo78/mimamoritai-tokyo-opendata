"""提出デッキに「誤報率」のスライドを1枚挿入する。

安全設計の3枚（7〜9枚目）は「ルールで止めている」という主張で終わっている。
その次に置くのは、その主張を数字で確かめた1枚。挿入位置は10枚目。
意匠は add-hardware-slide.py と同じ（Berry & Cream / Georgia見出し・Meiryo本文）。

数字は docs/eval/false-alarm-rate.md の実行結果と一致させること。
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DECK = "docs/mimamoritai-deck.pptx"
INSERT_AT = 9  # 0-based。安全設計3枚の直後＝10枚目になる

INK = RGBColor(0x2B, 0x2B, 0x2B)
BERRY = RGBColor(0x6D, 0x2E, 0x46)
ROSE = RGBColor(0xA2, 0x67, 0x69)
MUTED = RGBColor(0x7A, 0x7A, 0x7A)
FAINT = RGBColor(0xB9, 0xB9, 0xB9)
CREAM = RGBColor(0xEC, 0xE2, 0xD0)
HAIR = RGBColor(0xE3, 0xE3, 0xE3)
PAPER = RGBColor(0xFA, 0xF8, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (s, size, bold, font, color) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return box


def block(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


# 表の座標。左パネル全体は x=0.85 から幅 6.35。
COL_X = [0.85, 3.40, 5.30]
COL_W = [2.55, 1.90, 1.90]
HDR_Y, HDR_H = 1.86, 0.44
ROW_Y = [2.30, 3.24]
ROW_H = 0.94


def matrix(s):
    """平常日と異常日を、通知したか否かで4つに割った表。"""
    # 見出し行
    for i, label in enumerate(["通知した", "通知しなかった"], start=1):
        block(s, COL_X[i], HDR_Y, COL_W[i], HDR_H, PAPER, HAIR)
        text(s, COL_X[i], HDR_Y, COL_W[i], HDR_H,
             [(label, 10.5, False, "Meiryo", MUTED)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ("平常日", "通知は不要 ・ 32 件", [("0", "誤報", True), ("32", "", False)]),
        ("異常日", "通知が必要 ・ 20 件", [("20", "", False), ("0", "見逃し", True)]),
    ]
    for r, (head, sub, cells) in enumerate(rows):
        y = ROW_Y[r]
        block(s, COL_X[0], y, COL_W[0], ROW_H, PAPER, HAIR)
        block(s, COL_X[0], y, 0.05, ROW_H, BERRY if r else ROSE)
        text(s, COL_X[0] + 0.24, y + 0.20, COL_W[0] - 0.44, 0.28,
             [(head, 13.5, True, "Meiryo", INK)])
        text(s, COL_X[0] + 0.24, y + 0.54, COL_W[0] - 0.44, 0.26,
             [(sub, 10, False, "Meiryo", MUTED)])

        for i, (num, tag, hero) in enumerate(cells, start=1):
            block(s, COL_X[i], y, COL_W[i], ROW_H, CREAM if hero else WHITE, HAIR)
            # 0 の側だけ大きく出す。達成はここにしかない。
            text(s, COL_X[i], y + (0.14 if tag else 0.24), COL_W[i], 0.52,
                 [(num, 30 if hero else 24, True, "Georgia", BERRY if hero else FAINT)],
                 align=PP_ALIGN.CENTER)
            if tag:
                text(s, COL_X[i], y + 0.66, COL_W[i], 0.24,
                     [(tag, 10, True, "Meiryo", ROSE)], align=PP_ALIGN.CENTER)


def hero(s):
    """右側。誤報率を主役にする。"""
    x, w = 7.55, 4.93
    top, bottom = HDR_Y, ROW_Y[1] + ROW_H
    block(s, x, top, w, bottom - top, CREAM)
    block(s, x, top, 0.06, bottom - top, BERRY)

    text(s, x + 0.42, top + 0.20, w - 0.84, 0.3,
         [("誤報率", 13, True, "Meiryo", ROSE)])
    text(s, x + 0.42, top + 0.48, w - 0.84, 1.0,
         [("0.0 %", 54, True, "Georgia", BERRY)])
    text(s, x + 0.42, top + 1.42, w - 0.84, 0.28,
         [("平常日 32 件すべてで沈黙。見逃しも 0 件", 11.5, True, "Meiryo", INK)])
    text(s, x + 0.42, top + 1.74, w - 0.84, 0.56,
         [("検知率は「全部に鳴らせば 100%」で作れる。", 10.5, False, "Meiryo", MUTED),
          ("難しいのは、何でもない日に黙っていること。", 10.5, False, "Meiryo", MUTED)])


def finding(s):
    """このスライドの主役。評価が自分のバグを見つけたこと。"""
    y, h = 4.52, 1.52
    block(s, 0.85, y, 11.63, h, PAPER, HAIR)
    block(s, 0.85, y, 0.06, h, BERRY)

    text(s, 1.20, y + 0.18, 11.0, 0.3,
         [("この評価は、最初の実行で自分のバグを 1 件見つけた", 15, True, "Meiryo", BERRY)])

    lines = [
        "「活動量が普段より少ない」の判定が、進行中の一日の集計を、終わった一日の平均と比べていた。",
        "朝は少なくて当たり前なので、全世帯が毎朝鳴る。しかも鳴る引き金は、起きて照明を点けたことだった。",
        "毎朝鳴る通知は読まれなくなり、読まれなくなった通知は、本当の日に効かない。",
    ]
    for i, line in enumerate(lines):
        text(s, 1.20, y + 0.56 + i * 0.30, 11.0, 0.28,
             [(line, 11.5, False, "Meiryo", INK)])


def build(prs):
    layout = next((l for l in prs.slide_layouts if l.name == "DEFAULT"), None)
    if layout is None:
        layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(layout)

    text(s, 0.85, 0.46, 11.63, 0.72,
         [("「鳴らさなかった日」を数えている", 31, True, "Georgia", INK)])
    text(s, 0.85, 1.18, 11.63, 0.34,
         [("見守りの寿命を決めるのは、検知率ではなく誤報率", 13.5, False, "Meiryo", ROSE)])

    matrix(s)
    hero(s)
    finding(s)

    block(s, 0.85, 6.20, 11.63, 0.52, CREAM)
    block(s, 0.85, 6.20, 0.06, 0.52, BERRY)
    text(s, 1.20, 6.20, 11.28, 0.52,
         [("判定は引数だけで完結する純関数。LLM も外部 API も使わないので、毎 push で CI が再計算する",
           12.5, True, "Meiryo", BERRY)],
         anchor=MSO_ANCHOR.MIDDLE)

    block(s, 0.85, 6.90, 11.63, 0.012, HAIR)
    text(s, 0.85, 6.98, 5.4, 0.3,
         [("見守り隊 / CareRoute AI", 9.5, False, "Meiryo", FAINT)])
    text(s, 11.48, 6.98, 1.0, 0.3,
         [("10", 9.5, False, "Meiryo", FAINT)], align=PP_ALIGN.RIGHT)
    return s


def move(prs, frm, to):
    ids = prs.slides._sldIdLst
    el = list(ids)[frm]
    ids.remove(el)
    ids.insert(to, el)


def renumber(prs):
    """右下に焼き込まれたページ番号を、実際の並び順に合わせる。"""
    fixed = 0
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame or sh.left is None:
                continue
            if sh.left < Inches(11.2) or sh.top < Inches(6.5):
                continue
            if sh.text_frame.text.strip().isdigit():
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text.strip().isdigit():
                            if r.text != str(i):
                                fixed += 1
                            r.text = str(i)
    return fixed


def main():
    prs = Presentation(DECK)
    before = len(prs.slides._sldIdLst)
    build(prs)
    move(prs, before, INSERT_AT)
    fixed = renumber(prs)
    prs.save(DECK)
    print(f"slides={len(prs.slides._sldIdLst)} renumbered={fixed}")


if __name__ == "__main__":
    sys.exit(main())
